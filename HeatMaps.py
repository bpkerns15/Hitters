import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
import six
import copy
import seaborn as sns
import matplotlib.pyplot as plt
from PIL import Image
import matplotlib.patches as patches
import os
import math
from scipy import interpolate
import matplotlib.lines as lines
from tkinter import filedialog
from scipy.spatial.qhull import QhullError
import warnings

warnings.filterwarnings("ignore")

opponent = 'SPR_DRI'

def get_csv():
    global opponent
    #csv_file = 'CSVs//OregonHitters.csv'
    csv_file = 'drifters2.csv'
    #csv_file =  filedialog.askopenfilename()
    csv_df = pd.read_csv(csv_file)
    csv_df = csv_df.drop(csv_df[csv_df.BatterTeam != opponent].index)
    return csv_df

def to_player_df(csv_df, name):
    player_df = csv_df[(csv_df['Batter'] == name)]
    return player_df

def player_df_to_side_df(player_df, side):
    side_df = player_df[(player_df['PitcherThrows'] == side)]
    return side_df

def data_frame_for_damage_chart(pitch_df):
    damage_df = pitch_df

    #Remove unnesassary columns
    remove_list = damage_df.columns.values.tolist()
    remove_list.remove("Batter")
    remove_list.remove("PlateLocHeight")
    remove_list.remove("PlateLocSide")
    remove_list.remove("BatterSide")
    remove_list.remove("ExitSpeed")
    damage_df = damage_df.drop(remove_list, axis=1)

    #Removes Rows where EV is nan
    drop_index = []
    for index, row in damage_df.iterrows():
        if math.isnan(row['ExitSpeed']) == True:
            drop_index.append(index)
    damage_df = damage_df.drop(drop_index)

    #Removes Rows where EV < 60
    damage_df = damage_df.drop(damage_df[damage_df.ExitSpeed < 55.0].index)

    #Switches from Pitcher view to catcher view
    #damage_df['PlateLocSide'] = (damage_df['PlateLocSide'] * -1)

    return damage_df

def damage_chart(damage_df, name, side):
    x = []
    y = []
    array = []
    #print(damage_df)
    for i in range(10):
        row = []
        for j in range(8):
            temp_df = damage_df
            top_limit = 4.267222-0.35819444*i
            bottom_limit = 4.267222-0.35819444*(i+1)
            left_limit = -1.4166667 + 0.35416667*j
            right_limit = -1.416667 + 0.35416667*(j+1)

            #drops data from df not in cell needed
            too_left = temp_df[(temp_df['PlateLocSide'] < left_limit)].index
            temp_df = temp_df.drop(too_left)
            too_right = temp_df[(temp_df['PlateLocSide'] > right_limit)].index
            temp_df = temp_df.drop(too_right)
            too_high = temp_df[(temp_df['PlateLocHeight'] > top_limit)].index
            temp_df = temp_df.drop(too_high)
            too_low = temp_df[(temp_df['PlateLocHeight'] < bottom_limit)].index
            temp_df = temp_df.drop(too_low)

            avg_ev_for_zone = round(temp_df["ExitSpeed"].mean(),1)
            
                

            row.append(avg_ev_for_zone)

        array.append(row)


    df = pd.DataFrame(array)
    #print(df)

    try:
        df.to_numpy()
        x = np.arange(0, df.shape[1])
        y = np.arange(0, df.shape[0])
        #mask invalid values
        df = np.ma.masked_invalid(df)
        xx, yy = np.meshgrid(x, y)
        #get only the valid values
        x1 = xx[~df.mask]
        y1 = yy[~df.mask]
        newarr = df[~df.mask]

        GD1 = interpolate.griddata((x1, y1), newarr.ravel(),
                                (xx, yy),
                                    method='linear', fill_value=55.0)
        df = pd.DataFrame(GD1)
    except (ValueError,QhullError) as e:
        df = pd.DataFrame(array)
        df = df.fillna(55.0)


    
    #print(df)

    fig, ax = plt.subplots()
    rect = patches.Rectangle((1.5, 1.5), 4, 6, linewidth=1, edgecolor='black', facecolor='none')
    ax.add_patch(rect)
    horizontal_line1 = lines.Line2D([1.5, 5.5], [5.5, 5.5], linewidth=1, color='black')
    horizontal_line2 = lines.Line2D([1.5, 5.5], [3.5, 3.5], linewidth=1, color='black')
    vertical_line1 = lines.Line2D([2.833333, 2.833333], [1.5, 7.5], linewidth=1, color='black')
    vertical_line2 = lines.Line2D([4.166667, 4.166667], [1.5, 7.5], linewidth=1, color='black')
    ax.add_line(horizontal_line1)
    ax.add_line(horizontal_line2)
    ax.add_line(vertical_line1)
    ax.add_line(vertical_line2)

    fig = plt.imshow(df, cmap = 'jet',vmin=55,vmax=95,  interpolation='bicubic')
    #cbar = plt.colorbar(fig)
    fig.axes.get_xaxis().set_visible(False)
    fig.axes.get_yaxis().set_visible(False)
    #plt.title("Exit Velo Heat Map")
    newpath = os.path.join("Plots", name)
    if not os.path.exists(newpath):
        os.makedirs(newpath)
    #saves plot in folder
    plt.axis('off')
    plt.savefig(os.path.join("Plots", name, side + 'HP_heatmap.png'), bbox_inches='tight', pad_inches = 0)

    # load the image
    img = Image.open(os.path.join("Plots", name, side + 'HP_heatmap.png'))

    # get the color you want to make transparent
    color_to_replace = (0, 0, 127)  # in this example, we want to replace all shades of blue

    # make the color transparent
    img = img.convert("RGBA")
    data = img.getdata()

    new_data = []
    for item in data:
        if item[0] == color_to_replace[0] and item[1] == color_to_replace[1] and item[2] == color_to_replace[2]:
            new_data.append((255, 255, 255, 0))  # replace the color with a fully transparent color
        else:
            new_data.append(item)

    img.putdata(new_data)

    # save the modified image
    img.save(os.path.join("Plots", name, side + 'HP_heatmap.png'))

    return

def make_presentation():
    prs = Presentation("Template.pptx")
    prs.slide_width = Inches(10.5)
    prs.slide_height = Inches(8.115)
    #slide = prs.slides.add_slide(prs.slide_layouts[5])

    return prs

def stat_calcs(player_df,type):
    
    pitches = player_df.shape[0]

    #Swing %
    swings = (player_df["PitchCall"] == "InPlay").sum() + (player_df["PitchCall"] == "FoulBall").sum() + (player_df["PitchCall"] == "StrikeSwinging").sum()
    takes = (player_df["PitchCall"] == "BallCalled").sum() + (player_df["PitchCall"] == "StrikeCalled").sum() + (player_df["PitchCall"] == "HitByPitch").sum() + (player_df["PitchCall"] == "BallinDirt").sum()
    swing_rate = "%.1f" % round(100*(swings/(swings+takes)),1)
    whiffs = (player_df["PitchCall"] == "StrikeSwinging").sum()

    whiff_rate = "%.0f" % round(100*(whiffs/swings),0)

    fp = player_df[(player_df['Strikes'] == 0) & (player_df['Balls'] == 0)]
    fpswings = (fp["PitchCall"] == "InPlay").sum() + (fp["PitchCall"] == "FoulBall").sum() + (fp["PitchCall"] == "StrikeSwinging").sum()
    fptakes = (fp["PitchCall"] == "BallCalled").sum() + (fp["PitchCall"] == "StrikeCalled").sum() + (fp["PitchCall"] == "HitByPitch").sum() + (fp["PitchCall"] == "BallinDirt").sum()
    fpswing_rate = "%.0f" % round(100*(fpswings/(fpswings+fptakes)),0)

    #Chase Rate
    out_of_zone_df = player_df
    indexzone  = out_of_zone_df[ (out_of_zone_df['PlateLocSide'] < 0.83083) & (out_of_zone_df['PlateLocSide'] > -0.83083) & (out_of_zone_df['PlateLocHeight'] < 3.67333) & (out_of_zone_df['PlateLocHeight'] > 1.52417)].index
    out_of_zone_df = out_of_zone_df.drop(indexzone)
    chases = (out_of_zone_df["PitchCall"] == "InPlay").sum() + (out_of_zone_df["PitchCall"] == "FoulBall").sum() + (out_of_zone_df["PitchCall"] == "StrikeSwinging").sum()
    takes_out_of_zone =  (out_of_zone_df["PitchCall"] == "BallCalled").sum() + (out_of_zone_df["PitchCall"] == "StrikeCalled").sum() + (out_of_zone_df["PitchCall"] == "HitByPitch").sum() + (out_of_zone_df["PitchCall"] == "BallinDirt").sum()
    chase_rate = "%.0f" % round(100*(chases/(chases+takes_out_of_zone)),0)

    hbps = (player_df["PitchCall"] == "HitByPitch").sum()

    walk_df = player_df.drop(player_df[player_df.PitchCall == 'HitByPitch'].index)
    strikeouts = (player_df["KorBB"] == "Strikeout").sum()
    plate_apearences = (walk_df["KorBB"] == "Walk").sum() + (player_df["KorBB"] == "Strikeout").sum() + (player_df["PitchCall"] == "InPlay").sum() + (player_df["PitchCall"] == "HitByPitch").sum()

    walks = (walk_df["KorBB"] == "Walk").sum()
    bb_over_ks = "%.2f" % round(((walks)/strikeouts),2)

    #AVG
    ABs = (player_df["PitchCall"] == "InPlay").sum() - (player_df["PlayResult"] == "Sacrifice").sum() + (player_df["KorBB"] == "Strikeout").sum()
    hits = (player_df["PlayResult"] == "Single").sum() + (player_df["PlayResult"] == "Double").sum() + (player_df["PlayResult"] == "Triple").sum() + (player_df["PlayResult"] == "HomeRun").sum()
    avg = "%.3f" % round(hits/ABs,3)

    #Doubles, Triples, Home Runs
    doubles = (player_df["PlayResult"] == "Double").sum()
    triples = (player_df["PlayResult"] == "Triple").sum()
    homers = (player_df["PlayResult"] == "HomeRun").sum()

    #SLG
    slg = "%.3f" % round((hits + doubles + triples*2 + homers*3)/ABs,3)

    #OPS
    ops = "%.3f" % (round((walks + hbps + hits)/plate_apearences,3) + round((hits + doubles + triples*2 + homers*3)/ABs,3))

    if type == 'pitch':
        return_list = [str(pitches),str(avg.lstrip('0')),str(slg.lstrip('0')),str(whiff_rate + '%'),str(chase_rate + '%'),str(fpswing_rate + '%')]
    if type == 'side':
        return_list = [str(plate_apearences),str(avg.lstrip('0')),str(slg.lstrip('0')),str(whiff_rate + '%'),str(chase_rate + '%'),str(fpswing_rate + '%'),str(ops.lstrip('0')),str(bb_over_ks)]
    
    return return_list

def ev_calculations(player_df, lhp_df, rhp_df):

    lhp_df = lhp_df[lhp_df['ExitSpeed'] >= 55.0]
    rhp_df = rhp_df[rhp_df['ExitSpeed'] >= 55.0]
    player_df = player_df[player_df['ExitSpeed'] >= 55.0]
    lhp_df.dropna(subset=['ExitSpeed'], inplace=True)
    rhp_df.dropna(subset=['ExitSpeed'], inplace=True)
    player_df.dropna(subset=['ExitSpeed'], inplace=True)
    

    max_exit_speed = str(round(player_df['ExitSpeed'].max(),1))
    average_exit_speed = str(round(player_df['ExitSpeed'].mean(),1))
    average_exit_speed_rhp = str(round(rhp_df['ExitSpeed'].mean(),1))
    average_exit_speed_lhp = str(round(lhp_df['ExitSpeed'].mean(),1))

    over_90_lhp = lhp_df[lhp_df['ExitSpeed'] >= 90.0]
    over_90_rhp = rhp_df[rhp_df['ExitSpeed'] >= 90.0]
    if rhp_df.shape[0] > 0:
        hard_hit_rhp = str(round((over_90_rhp.shape[0]/rhp_df.shape[0])*100,1))
    else:
        hard_hit_rhp = 'N/A'

    if lhp_df.shape[0] > 0:
        hard_hit_lhp = str(round((over_90_lhp.shape[0]/lhp_df.shape[0])*100,1))
    else:
        hard_hit_lhp = 'N/A'

    return_list = [max_exit_speed,average_exit_speed_lhp,average_exit_speed,average_exit_speed_rhp,hard_hit_lhp,hard_hit_rhp]

    return return_list

def ground_out_fly_out(df):
    # Filter out rows where ExitSpeed is below 50 and PlayResult is not 'out'
    filtered_df = df[(df['ExitSpeed'] >= 50) & (df['PlayResult'] == 'Out')]

    # Count the number of rows where Angle is less than 10
    angle_less_than_10 = filtered_df[filtered_df['Angle'] < 10].shape[0]

    # Count the number of rows where Angle is greater than 10
    angle_greater_than_10 = filtered_df[filtered_df['Angle'] > 10].shape[0]

    # Compute the ratio
    if angle_greater_than_10 != 0:
        ratio = angle_less_than_10 / angle_greater_than_10
        ratio = str("%.2f" % round(ratio,2))
    else:
        ratio = str(str(angle_less_than_10) + '/0')

    return ratio

def up_in_zone_avg(df):
    upzone = df[(df['PlateLocSide'] < 0.83083) & (df['PlateLocSide'] > -0.83083) & (df['PlateLocHeight'] < 3.67333) & (df['PlateLocHeight'] > 2.95694)]
    #AVG
    ABs = (upzone["PitchCall"] == "InPlay").sum() - (upzone["PlayResult"] == "Sacrifice").sum() + (upzone["KorBB"] == "Strikeout").sum()
    hits = (upzone["PlayResult"] == "Single").sum() + (upzone["PlayResult"] == "Double").sum() + (upzone["PlayResult"] == "Triple").sum() + (upzone["PlayResult"] == "HomeRun").sum()
    avg = "%.3f" % round(hits/ABs,3)

    return str(avg)

def down_in_zone_avg(df):
    down_zone = df[(df['PlateLocSide'] < 0.83083) & (df['PlateLocSide'] > -0.83083) & (df['PlateLocHeight'] < 2.24056) & (df['PlateLocHeight'] > 1.52417)]
    #AVG
    ABs = (down_zone["PitchCall"] == "InPlay").sum() - (down_zone["PlayResult"] == "Sacrifice").sum() + (down_zone["KorBB"] == "Strikeout").sum()
    hits = (down_zone["PlayResult"] == "Single").sum() + (down_zone["PlayResult"] == "Double").sum() + (down_zone["PlayResult"] == "Triple").sum() + (down_zone["PlayResult"] == "HomeRun").sum()
    avg = "%.3f" % round(hits/ABs,3)

    return str(avg)

def breaking_ball_df(side_df):
    bb_df = side_df[(side_df['TaggedPitchType'].isin(['BreakingBall','Slider','Curveball']))]
    return bb_df

def change_up_df(side_df):
    ch_df = side_df[(side_df['TaggedPitchType'].isin(['ChangeUp','Splitter','Knuckleball','OffSpeed','Changeup']))]
    return ch_df

def fastball_df(side_df):
    fb_df = side_df[(side_df['TaggedPitchType'].isin(['Fastball','Sinker','OneSeamFastBall','TwoSeamFastBall','FourSeamFastBall', 'Cutter']))]
    #print(fb_df)
    return fb_df

def presentation(name, prs, ev_stats,rhp_stats,lhp_stats,rhp_fb_stats,lhp_fb_stats,rhp_bb_stats,lhp_bb_stats,rhp_ch_stats,lhp_ch_stats,overall_stats,go_fo, i):

    ops_stats = [lhp_stats[6],overall_stats[6],rhp_stats[6]]
    k_bb = [lhp_stats[7],overall_stats[7],rhp_stats[7]]


    first_slide_id = prs.slides[0].slide_id
    slide = prs.slides.get(first_slide_id+i)
    name1 = name.split()
    full_name = name1[-1] + ' ' + name1[0]
    full_name = full_name[:-1]
    title = slide.shapes.title
    title.text = full_name
    title.text_frame.paragraphs[0].font.size = Pt(17)
    title.text_frame.paragraphs[0].font.bold = False
    title.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 248, 221)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.vertical_anchor = MSO_ANCHOR.MIDDLE

    rhp = os.path.join("Plots", name, 'Right' + 'HP_heatmap.png')
    lhp = os.path.join("Plots", name, 'Left' + 'HP_heatmap.png')
    shape = slide.shapes.add_picture(lhp,Inches(1.94),Inches(3.48), width=Inches(2.2), height=Inches(2.76))
    shape = slide.shapes.add_picture(rhp,Inches(6.37),Inches(3.48), width=Inches(2.2), height=Inches(2.76))

    #LHP top table
    x, y, cx, cy = Inches(0.01), Inches(0.84), Inches(4.04), Inches(2.55)
    shape = slide.shapes.add_table(5, 6, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    tbl[0][-1].text = style_id
    text_labels = ['LHP\n' + lhp_stats[0] + ' PA','AVG','SLG','Whiff %', 'Chase %', 'FPS %']
    for i in range(6):
        cell = table.cell(0, i)
        cell.text = text_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        if i == 3 or 4:
            cell.text_frame.paragraphs[0].font.size = Pt(12.5)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(152, 1, 46)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 248, 221)
        if i == 0:
            cell.fill.fore_color.rgb = RGBColor(255, 248, 221)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(94, 94, 94)

    pitch_labels = ['FB\n' + '(' + lhp_fb_stats[0] + ')','BB\n' + '(' + lhp_bb_stats[0] + ')', 'CH\n' + '(' + lhp_ch_stats[0] + ')']
    for i in range(3):
        cell = table.cell(i+1, 0)
        cell.text = pitch_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 248, 221)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(94, 94, 94)
    
    cell = table.cell(4, 0)
    cell.text = 'LHP Overall'
    cell.text_frame.paragraphs[0].font.size = Pt(9)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(94, 94, 94)

    lhp_pitch_values = [lhp_fb_stats,lhp_bb_stats,lhp_ch_stats]
    for i in range(3):
        stats = lhp_pitch_values[i]
        for j in range(5):
            cell = table.cell(i+1, j+1)
            cell.text = stats[j+1]
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    for i in range(5):
        cell = table.cell(4, i+1)
        cell.text = lhp_stats[i+1]
        cell.text_frame.paragraphs[0].font.size = Pt(12.5)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(201, 201, 201)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            

    #RHP top table
    x, y, cx, cy = Inches(6.45), Inches(0.84), Inches(4.04), Inches(2.55)
    shape = slide.shapes.add_table(5, 6, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    tbl[0][-1].text = style_id
    text_labels = ['RHP\n' + rhp_stats[0] + ' PA','AVG','SLG','Whiff %', 'Chase %', 'FPS %']
    for i in range(6):
        cell = table.cell(0, i)
        cell.text = text_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        if i == 3 or 4:
            cell.text_frame.paragraphs[0].font.size = Pt(12.5)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(152, 1, 46)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 248, 221)
        if i == 0:
            cell.fill.fore_color.rgb = RGBColor(255, 248, 221)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(94, 94, 94)

    pitch_labels = ['FB\n' + '(' + rhp_fb_stats[0] + ')','BB\n' + '(' + rhp_bb_stats[0] + ')', 'CH\n' + '(' + rhp_ch_stats[0] + ')']
    for i in range(3):
        cell = table.cell(i+1, 0)
        cell.text = pitch_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 248, 221)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(94, 94, 94)
    
    cell = table.cell(4, 0)
    cell.text = 'RHP Overall'
    cell.text_frame.paragraphs[0].font.size = Pt(9)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(94, 94, 94)

    rhp_pitch_values = [rhp_fb_stats,rhp_bb_stats,rhp_ch_stats]
    for i in range(3):
        stats = rhp_pitch_values[i]
        for j in range(5):
            cell = table.cell(i+1, j+1)
            cell.text = stats[j+1]
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    for i in range(5):
        cell = table.cell(4, i+1)
        cell.text = rhp_stats[i+1]
        cell.text_frame.paragraphs[0].font.size = Pt(12.5)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(201, 201, 201)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    #EVs RHP Hard Hit%
    x, y, cx, cy = Inches(6.41), Inches(6.73), Inches(1.32), Inches(0.99)
    shape = slide.shapes.add_table((2), 1, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tbl[0][-1].text = style_id

    cell = table.cell(0, 0)
    cell.text = 'Hard Hit % RHP'
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    cell = table.cell(1, 0)
    cell.text = ev_stats[5]
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    #LHP Hard Hit%
    x, y, cx, cy = Inches(2.79), Inches(6.73), Inches(1.32), Inches(0.99)
    shape = slide.shapes.add_table((2), 1, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    tbl[0][-1].text = style_id

    cell = table.cell(0, 0)
    cell.text = 'Hard Hit % LHP'
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    cell = table.cell(1, 0)
    cell.text = ev_stats[4]
    cell.text_frame.paragraphs[0].font.size = Pt(14)
    cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    #top table
    x, y, cx, cy = Inches(4.2), Inches(2.33), Inches(2.14), Inches(1.06)
    shape = slide.shapes.add_table((2), 3, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    tbl[0][-1].text = style_id

    text_labels = ['SB/ATT','OPS Last 10 G','Max EV']
    bottom_labels = ['','',ev_stats[0]]
    for i in range(3):
        cell = table.cell(0, i)
        cell.text = text_labels[i]
        if i == 0:
            size = 11
        elif i == 1:
            size = 9
        elif i == 2:
            size = 14
        cell.text_frame.paragraphs[0].font.size = Pt(size)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(152, 1, 46)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 248, 221)

        cell = table.cell(1, i)
        cell.text = bottom_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    #ev avgs
    x, y, cx, cy = Inches(4.2), Inches(6.92), Inches(2.14), Inches(1.06)
    shape = slide.shapes.add_table((2), 3, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    tbl[0][-1].text = style_id

    text_labels = ['LHP EV','AVG EV','RHP EV']
    for i in range(3):
        cell = table.cell(0, i)
        cell.text = text_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(152, 1, 46)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 248, 221)

        cell = table.cell(1, i)
        cell.text = ev_stats[1+i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    #OPS
    x, y, cx, cy = Inches(4.2), Inches(3.48), Inches(2.14), Inches(1.06)
    shape = slide.shapes.add_table((2), 3, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    tbl[0][-1].text = style_id

    text_labels = ['OPS LHP','OPS Overall','OPS RHP']
    for i in range(3):
        cell = table.cell(0, i)
        cell.text = text_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        if i == 1:
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(152, 1, 46)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 248, 221)

        cell = table.cell(1, i)
        cell.text = ops_stats[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    #K/BB
    x, y, cx, cy = Inches(4.2), Inches(4.64), Inches(2.14), Inches(1.06)
    shape = slide.shapes.add_table((2), 3, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    tbl[0][-1].text = style_id

    text_labels = ['K/BB LHP','K/BB Overall','K/BB RHP']
    for i in range(3):
        cell = table.cell(0, i)
        cell.text = text_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        if i == 1:
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(152, 1, 46)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 248, 221)

        cell = table.cell(1, i)
        cell.text = k_bb[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    #GO/FO
    x, y, cx, cy = Inches(4.2), Inches(5.81), Inches(2.14), Inches(1.06)
    shape = slide.shapes.add_table((2), 3, x, y, cx, cy)
    table = shape.table
    tbl =  shape._element.graphic.graphicData.tbl
    style_id = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    tbl[0][-1].text = style_id

    text_labels = ['GO/FO LHP','GO/FO Overall','GO/FO RHP']
    for i in range(3):
        cell = table.cell(0, i)
        cell.text = text_labels[i]
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        if i == 1:
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = False
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(152, 1, 46)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 248, 221)

        cell = table.cell(1, i)
        cell.text = go_fo[i]
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.name = 'Adobe Heiti Std R'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)


    newpath = os.path.join("FinalPres", opponent)
    if not os.path.exists(newpath):
        os.makedirs(newpath)

    prs.save(os.path.join("FinalPres", opponent, opponent + '.pptx'))

    return prs

def main():
    csv = get_csv()
    prs = make_presentation()
    names = csv["Batter"].unique()
    for i in range(len(names)):
        print(names[i])
        player_df = to_player_df(csv,names[i])
        rhp_df = player_df_to_side_df(player_df,'Right')
        rhp_fb_df = fastball_df(rhp_df)
        rhp_bb_df = breaking_ball_df(rhp_df)
        rhp_ch_df = change_up_df(rhp_df)
        lhp_df = player_df_to_side_df(player_df,'Left')
        lhp_fb_df = fastball_df(lhp_df)
        lhp_bb_df = breaking_ball_df(lhp_df)
        lhp_ch_df = change_up_df(lhp_df)
        ev_stats = ev_calculations(player_df,lhp_df,rhp_df)
        lhp_df_dc = data_frame_for_damage_chart(lhp_df)
        rhp_df_dc = data_frame_for_damage_chart(rhp_df)
        damage_chart(lhp_df_dc,names[i],'Left')
        damage_chart(rhp_df_dc,names[i],'Right')
        print('LHP UP:' + up_in_zone_avg(lhp_df))
        print('LHP Down:' + down_in_zone_avg(lhp_df))
        print('RHP UP:' + up_in_zone_avg(rhp_df))
        print('RHP Down:' + down_in_zone_avg(rhp_df))
        prs = presentation(names[i],prs,ev_stats,rhp_stats=stat_calcs(rhp_df,'side'),lhp_stats=stat_calcs(lhp_df,'side'),rhp_fb_stats=stat_calcs(rhp_fb_df,'pitch'),lhp_fb_stats=stat_calcs(lhp_fb_df,'pitch'),rhp_bb_stats=stat_calcs(rhp_bb_df,'pitch'),lhp_bb_stats=stat_calcs(lhp_bb_df,'pitch'),rhp_ch_stats=stat_calcs(rhp_ch_df,'pitch'),lhp_ch_stats=stat_calcs(lhp_ch_df,'pitch'),overall_stats=stat_calcs(player_df,'side'),go_fo=[ground_out_fly_out(lhp_df),ground_out_fly_out(player_df),ground_out_fly_out(rhp_df)],i=i)
        
main()




